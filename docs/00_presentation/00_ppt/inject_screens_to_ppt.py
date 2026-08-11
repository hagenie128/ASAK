# -*- coding: utf-8 -*-
"""스크린샷을 수정2 PPT 복사본에 삽입한다."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent
SHOT = ROOT / "screenshots"
SRC = ROOT / "ASAK_샐러드_스마트키오스크_수정2_20260811.pptx"
DST = ROOT / "ASAK_샐러드_스마트키오스크_수정3_20260811.pptx"

# (slide 1-based, shape_index, image_filename)
# shape_index는 _shapes.txt 기준 자리표시 사각형
PLACEMENTS: list[tuple[int, int, str]] = [
    # 11 화면설계서 — 3프레임
    (11, 23, "kiosk-02-menuList.png"),
    (11, 25, "kiosk-03-menuDetail.png"),
    (11, 27, "kiosk-04-cart.png"),
    # 12 메뉴 목록
    (12, 19, "kiosk-02-menuList.png"),
    (12, 22, "kiosk-02-menuList.png"),
    # 13 메뉴 상세
    (13, 19, "kiosk-03-menuDetail.png"),
    (13, 22, "kiosk-03-menuDetail.png"),
    # 14 장바구니·결제
    (14, 19, "kiosk-04-cart.png"),
    (14, 22, "kiosk-05-payment.png"),
    # 15 Admin 주문
    (15, 19, "admin-01-liveOrders.png"),
    # 16 Admin 메뉴
    (16, 19, "admin-03-menus.png"),
    # 19 결과물 4칸 (안쪽 프레임)
    (19, 13, "kiosk-02-menuList.png"),
    (19, 18, "kiosk-03-menuDetail.png"),
    (19, 23, "admin-01-liveOrders.png"),
    (19, 28, "admin-03-menus.png"),
]

# 자리표시 라벨 텍스트를 비울 슬라이드·인덱스
CLEAR_TEXT: dict[int, list[int]] = {
    11: [],  # 프레임에 라벨 없음
    12: [20, 21, 23, 24],
    13: [20, 21, 23, 24],
    14: [20, 21, 23, 24],
    15: [20, 21],
    16: [20, 21],
    19: [14, 15, 19, 20, 24, 25, 29, 30],
}


def clear_shape_text(shape) -> None:
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.text = ""
        if not p.runs:
            p.text = ""


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"missing {SRC}")
    for _, _, name in PLACEMENTS:
        if not (SHOT / name).exists():
            raise SystemExit(f"missing shot {name}")

    shutil.copy2(SRC, DST)
    prs = Presentation(str(DST))

    for slide_no, shape_idx, fname in PLACEMENTS:
        slide = prs.slides[slide_no - 1]
        shapes = list(slide.shapes)
        if shape_idx >= len(shapes):
            print(f"SKIP S{slide_no}[{shape_idx}] out of range")
            continue
        box = shapes[shape_idx]
        left, top, width, height = box.left, box.top, box.width, box.height
        img = SHOT / fname
        slide.shapes.add_picture(str(img), Emu(left), Emu(top), width=Emu(width), height=Emu(height))
        print(f"S{slide_no}[{shape_idx}] ← {fname}")

    for slide_no, idxs in CLEAR_TEXT.items():
        slide = prs.slides[slide_no - 1]
        shapes = list(slide.shapes)
        for idx in idxs:
            if idx < len(shapes):
                clear_shape_text(shapes[idx])

    prs.save(str(DST))
    print(f"saved {DST}")


if __name__ == "__main__":
    main()
