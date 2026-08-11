# -*- coding: utf-8 -*-
"""수정1 PPTX 복사본에 ASAK_PPT_작성계획_초안.md 문구를 반영한다.
레이아웃·이미지·서식은 유지하고, 텍스트 박스 전체 문자열만 교체한다.
"""
from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC = Path(__file__).with_name("ASAK_샐러드_스마트키오스크_수정1_20260811.pptx")
DST = Path(__file__).with_name("ASAK_샐러드_스마트키오스크_수정2_20260811.pptx")

# (슬라이드 1-based, old_full_text, new_full_text)
# old는 dump 기준 정확한 전체 문자열. \x0b는 수직탭(문단 구분).
REPLACEMENTS: list[tuple[int, str, str]] = [
    # ── 1 표지 ──
    (
        1,
        "옵션 선택이 많은 샐러드 주문을 고객 키오스크와 관리자 운영 화면으로 연결한다",
        "A Salad A Kiosk — 옵션이 많은 샐러드 주문을, 키오스크와 운영 화면으로 연결한다",
    ),
    # ── 2 목차 부연 ──
    (2, "주제 · 배경 · 범위 · 기대 효과", "Why — 왜 샐러드 키오스크인가"),
    (2, "담당 화면 · API · 산출물", "Who — 화면 · API 단위 역할"),
    (2, "기획 → 설계 → 개발 → 통합 → QA", "How — 문서 → 코드 → 검증"),
    (2, "기술 · 구조 · ERD · 화면 · 결과물", "Build — 구조 · 화면 · 결과물"),
    (2, "잘한 점 · 보완할 점 · 배운 점", "Review — 정직한 현재 단계"),
    # ── 3 개요 ──
    (
        3,
        "옵션 선택이 많은 샐러드 주문을 고객 키오스크와 관리자 운영 화면으로 연결한다.  ·  2인 팀 · 6주 · 저장소 4개",
        "옵션이 많은 샐러드 주문을 실수 없이 끝낸다.  ·  2인 팀 · 6주 · 저장소 4개",
    ),
    (3, "주제 및\n선정 배경", "주제 ·\n기획 의도"),
    (
        3,
        "샐러드 주문은 메뉴 외에도\n옵션 선택, 재료 제외,\n수량 변경처럼\n확인할 항목이 많다.\x0b\x0b주문 과정을 쉽게\n이해할 수 있는 흐름을\n만들고자 했다.\n",
        "샐러드 주문은 메뉴 외에도\n옵션·제외·수량이\n겹친다.\x0b\x0b실수 없이 끝내는\n주문 흐름이\n필요했다.",
    ),
    (3, "특화 포인트\n차별화 내용", "특화 ·\n차별화"),
    (
        3,
        "단순 메뉴 주문이 아니라\n옵션 정책, 제외 재료,\n품절 여부, 수량 제한을\n주문 흐름에 반영한다.\x0b\x0b주문 결과는 관리자\n화면에서 확인한다.",
        "옵션 정책·제외 재료·\n품절·수량 제한을\n주문 흐름과\n서버 검증에 반영한다.\x0b\x0b결과는 Admin에서\n확인한다.",
    ),
    (3, "Home · 메뉴 목록\n메뉴 상세 · 장바구니 결제", "Home · 목록 · 상세\n장바구니 · 결제"),
    (3, "메뉴 · 옵션 · 주문\n데이터 검증", "메뉴 · 옵션 · 주문\n검증 API"),
    (3, "활용 툴 및 자료", "활용 기술 ·\n자료"),
    (
        3,
        "고객은 복잡한 옵션을\n화면 흐름에 따라\n선택할 수 있다.",
        "고객: 옵션을\n흐름대로 선택한다.",
    ),
    (
        3,
        "운영자는 주문과 메뉴를\n관리자 화면에서\n확인할 수 있다.",
        "운영: 주문·메뉴를\n화면에서 확인한다.",
    ),
    (
        3,
        "주문 입력 오류를\n줄이는 것이 목표다.",
        "목표: 입력 오류를\n줄인다.",
    ),
    # ── 4 팀 ──
    (
        4,
        "2인 팀으로 저장소 4개(문서 · Kiosk · Admin · Backend)를 나누어 맡고, WBS 85건을 증거 기준으로 관리했다.",
        "2인 팀 · 저장소 4개(ASAK · Kiosk · Admin · back) · 역할은 화면·API 단위로 나눈다.",
    ),
    (4, "기획 · 문서 · 관리자 화면", "Admin · Docs"),
    (4, "고객 화면 · 백엔드", "Kiosk · QA"),
    (
        4,
        "프로젝트 문서·WBS·ERD·API 명세 정본 관리\nFigma 디자인 토큰 · 공통 컴포넌트 · 화면 매핑\nAdmin 전 화면 구현 — Live 주문 · 주문 관리 · 상태 전환\n품절 · 메뉴 · 결제수단 · 매출 · 대시보드 화면\n관리자 실서버 API 연동, DB 스키마 반영",
        "Admin 화면 · 문서 · Hub 정본 관리\nFigma · Screen Bible · WBS·발표 자료\nLive 주문 · 주문 관리 · 메뉴·품절 화면\n관리자 API 연동 · DB 스키마 반영",
    ),
    (
        4,
        "관리자 화면 Figma 초기 설계 (로그인 · 주문 현황 · 매출)\nKiosk 전 화면 구현 — 라우트 · 메뉴 목록 · 상세 · 옵션\n필수·선택 옵션 활성 조건\n수량 제한(메뉴 9개 · 장바구니 30개) 및 초과 안내\n결제 · 주문 완료 · 타임아웃 화면, Kiosk 실서버 API 연동",
        "Kiosk 주문 흐름 · 옵션 UI\n필수·선택 옵션 · 제외 · 수량 제한\n장바구니 · 결제·완료·타임아웃 화면\nAPI 연동·검증 · QA·시연 시나리오",
    ),
    (
        4,
        "멘토 남상규  ·  주차별 강사 피드백으로 API 계약 · 상태 관리 · 수직 슬라이스 구현 순서를 조정",
        "멘토 남상규  ·  강사 피드백으로 API 계약 · 상태 경계 · 수직 슬라이스 순서를 맞췄다",
    ),
    # ── 5 절차 ──
    (
        5,
        "주차별 회의와 wbs를 활용한 체계적인 업무 분담\x0b7단계 실서버 연동과 8단계 통합 검증은 현재 진행중",
        "문서(요구·화면·API)를 먼저 맞춘 뒤 화면과 서버를 붙였다.\x0b7 실서버 연동 · 8 통합 검증은 현재 진행 중",
    ),
    # ── 6 경과 요약 ──
    (
        6,
        "ASAK 프로젝트 진행 핵심 사항",
        "Stack · Build · Process · Iterate · Deliver — 다섯 축으로 경과를 본다.",
    ),
    (
        6,
        "React 기반 Kiosk /\nAdmin 화면과\nSpring Boot · MyBatis ·\nMySQL 기반 API 구조 사용",
        "React Kiosk/Admin +\nSpring · MyBatis ·\nMySQL API 구조",
    ),
    (
        6,
        "메뉴 · 옵션 ·\n제외 재료 · 수량 선택,\n장바구니,\n관리자 주문 · 메뉴\n관리 화면 구성",
        "옵션·제외·수량 ·\n장바구니 ·\nAdmin 주문·메뉴",
    ),
    (
        6,
        "요구사항 분석부터\nUI/DB 설계, 화면 · API\n개발, 통합, QA,\n시연 준비까지\n단계별로 진행했다.",
        "요구 → 설계 →\n개발 → 통합 →\nQA → 시연",
    ),
    (
        6,
        "검토에서 확인된\n화면 흐름 · 주문 검증 ·\n데이터 구조 사항을\n보완했다.",
        "화면 흐름·검증·\n데이터 구조를\n피드백 단위로 보완",
    ),
    (
        6,
        "Kiosk, Admin, ERD,\nAPI 테스트,\nFigma · WBS 자료를\n결과물로 정리했다.",
        "화면 · ERD ·\nAPI 테스트 ·\nBible · WBS",
    ),
    # ── 7 기술 ──
    (
        7,
        "Home, 메뉴 목록 · 상세, 장바구니, 결제 관련 화면 및 주문 상태 관리",
        "주문 세션 · 목록·상세·카트·결제 UI",
    ),
    (
        7,
        "Live 주문, 주문 관리, 메뉴 관리, 품절 · 판매 관련 화면",
        "Live · 주문 · 메뉴 · 품절·판매 화면",
    ),
    (
        7,
        "키오스크 · 관리자 API, 메뉴 · 옵션 · 주문 검증 및 상태 처리",
        "메뉴·옵션·주문 검증 · Admin 조회 API",
    ),
    (
        7,
        "24개 테이블 — 메뉴 · 재료 · 옵션 정책 · 주문 · 결제, 상태값은 공통코드로 관리",
        "24테이블 — 메뉴·옵션·주문·결제 (공통코드)",
    ),
    (
        7,
        "형상 관리, 화면 설계, API 테스트, 작업 상태 관리",
        "설계 · 형상 · API 테스트 · 일정",
    ),
    (
        7,
        "Kiosk와 Admin은 역할별로 독립된 React 애플리케이션으로 구성하고, Spring Boot API를 통해 메뉴 · 옵션 · 주문 데이터를 조회 · 검증한다.\nMyBatis Mapper는 메뉴 구성과 주문 저장에 필요한 테이블 관계를 연결한다.",
        "Kiosk와 Admin은 역할별 React 앱으로 분리하고, Spring Boot API로 메뉴·옵션·주문을 조회·검증한다.\nMyBatis는 메뉴 구성과 주문 보존에 필요한 관계를 연결한다.",
    ),
    # ── 8 구조 ──
    (
        8,
        "Home · Menu · Detail · Cart · Payment\nReact Router · Zustand · Axios",
        "Home · Menu · Detail · Cart · Payment\n고객 주문 세션",
    ),
    (
        8,
        "Live 주문 · 주문 관리 · 메뉴 관리\n품절 · 결제수단 · 매출",
        "Live · Orders · Menus\n품절 · 결제수단 · 매출",
    ),
    # ── 9~10 ERD 캡션 ──
    (
        9,
        "메뉴는 카테고리 · 기본 재료(menu_ing) · 옵션 정책(menu_opt_policy)으로 구성하며, 옵션 선택 조건(min/max_select)과 제외 가능 재료(can_remove)를 주문 화면에 제공한다.",
        "메뉴 = 카테고리 + 기본 재료 + 옵션 정책.  min/max 선택과 can_remove가 주문 화면의 규칙이 된다.",
    ),
    (
        10,
        "주문은 메뉴 단위의 주문 항목으로 저장하고, 선택 옵션(order_item_option)과 제외 재료(item_exclusion)를 주문 항목에 연결해 주문 내용을 보존한다.",
        "주문 항목에 옵션·제외 재료를 붙여 ‘그때 그 조합’을 그대로 보존한다.",
    ),
    # ── 11 화면설계 ──
    (11, "카테고리 선택에서 장바구니까지", "키오스크 주문 한 줄 지도"),
    (
        11,
        "고객은 카테고리별 메뉴를 조회하고,\n상세 화면에서 옵션 · 제외 재료 ·\n수량을 선택합니다.",
        "카테고리 → 메뉴 →\n옵션·제외·수량 →\n장바구니 → 결제",
    ),
    (
        11,
        "선택 결과는 장바구니에 반영되어\n주문 금액을 확인할 수 있습니다.",
        "단계를 나누되, 뒤로 가도\n선택은 유지한다.",
    ),
    (
        11,
        "화면 상태(기본 · 품절 · 선택 불가)는 Screen Bible 기준으로 정의했습니다.",
        "칩 기준: 카테고리 · 옵션 정책 · 제외 재료 · 금액 확인  ·  Screen Bible",
    ),
    # ── 12 메뉴 목록 ──
    (
        12,
        "카테고리별로 메뉴를 조회하고, 각 메뉴의\n이름 · 가격 · 이미지 · 품절 여부를 확인합니다.",
        "카테고리별 메뉴를 조회하고\n이름 · 가격 · 이미지 · 품절을\n한눈에 확인합니다.",
    ),
    (
        12,
        "고객은 원하는 메뉴를 선택해\n상세 화면으로 이동합니다.",
        "원하는 메뉴를 선택해\n상세로 이동합니다.",
    ),
    (
        12,
        "메뉴 목록은 샐러드 주문 흐름의 시작점이며,\n품절 메뉴는 선택하지 못하도록 안내합니다.",
        "품절 메뉴는 선택 전에 막는다 — 주문 실패를 앞당기지 않는다.",
    ),
    # ── 13 상세·옵션 ──
    (
        13,
        "메뉴 기본 정보와 가격을 확인한 뒤 옵션 정책에\n따라 필수 · 선택 항목을 고릅니다.",
        "필수 · 선택 옵션을 정책(min/max)에\n맞게 고릅니다.",
    ),
    (
        13,
        "기본 재료 중 제외 가능한 재료를 선택하고,\n주문 수량을 변경할 수 있습니다.",
        "제외 가능한 재료만 빼고,\n수량을 조정합니다.\n가격 변화는 선택과 함께 확인합니다.",
    ),
    (
        13,
        "옵션 선택 조건 · 품절 여부 · 수량 제한은\n주문 생성 전 서버에서 다시 검증합니다.",
        "화면 규칙 = 서버 재검증. 옵션·품절·수량을 주문 전에 다시 본다.",
    ),
    # ── 14 장바구니 (시연 B: 결제 이동까지) ──
    (
        14,
        "장바구니에서 선택한 메뉴, 옵션, 제외 재료,\n수량, 주문 금액을 확인합니다.",
        "선택한 메뉴 · 옵션 · 제외 · 수량 ·\n금액을 확인합니다.",
    ),
    (
        14,
        "고객은 주문 내용을 검토한 뒤\n결제 화면으로 이동합니다.",
        "주문 내용을 검토한 뒤\n결제 화면으로 이동합니다.",
    ),
    (
        14,
        "결제 완료 처리 범위는 팀 확인 후 확정합니다.\n현재는 결제 화면 이동까지 표현합니다.",
        "장바구니는 마지막 검산. 실패해도 선택은 유지한다. (시연: 결제 화면 이동까지)",
    ),
    # ── 15 Admin 주문 ──
    (
        15,
        "관리자는 Live 주문과 주문 목록에서\n주문 번호, 주문 메뉴, 주문 상태를 확인합니다.",
        "Live 보드와 주문 목록에서\n번호 · 메뉴 · 상태를 확인합니다.",
    ),
    (
        15,
        "주문 상세에서는 선택 옵션과\n제외 재료 정보를 확인합니다.",
        "상세에서 옵션 · 제외 재료까지\n봅니다.",
    ),
    (
        15,
        "주문 상태는 접수 · 준비 · 완료\n흐름에 맞춰 관리합니다.",
        "고객 주문 → 운영 화면. 키오스크는 시작, Admin은 처리의 시작.",
    ),
    # ── 16 Admin 메뉴 ──
    (
        16,
        "관리자는 메뉴 목록 · 상세 · 카테고리\n정보를 조회합니다.",
        "메뉴 목록 · 상세 · 카테고리를\n조회합니다.",
    ),
    (
        16,
        "메뉴의 기본 정보와 옵션 정책은\n고객 주문 화면에 제공되는 데이터의 기준이 됩니다.",
        "옵션 정책은 고객 주문 화면\n데이터의 기준이 됩니다.",
    ),
    (
        16,
        "등록 · 수정 · 삭제 가능 범위는\n팀 확인 후 제목과 본문을 조정합니다.",
        "기준 데이터가 흔들리면 Kiosk도 흔들린다. (CRUD 범위는 구현 현황에 맞춤)",
    ),
    # ── 17 프로세스 하단 ──
    (
        17,
        "요구사항 분석  →  UI/DB 설계  →  키오스크 · 관리자 · 백엔드 개발  →  1차 통합  →  DB 연동  →  테스트 QA  →  발표 시연",
        "요구 → 설계 → 개발 → 통합 → DB → QA → 시연",
    ),
    # ── 19 시연 흐름 ──
    (
        19,
        "메뉴 선택 → 옵션 · 제외 재료 · 수량 선택 →\n장바구니 확인 → 주문 생성 → 관리자 주문 확인",
        "메뉴 → 옵션·제외·수량 → 장바구니 → (주문) → Admin 확인",
    ),
    # ── 20 검증 ──
    (
        20,
        "화면 설계, 데이터 모델, API 테스트 자료를 함께\n관리해 구현 화면과 데이터 흐름을 추적했습니다.",
        "설계 · 데이터 · 테스트 자료를 같이 남겨\n화면과 흐름을 추적했다.",
    ),
    # ── 21 자체 평가 ──
    (21, "실서버 연동\n진행 중", "통합 검증\n진행 중"),
    (
        21,
        "2인 팀이 문서 · Kiosk · Admin · Backend를 저장소로 분리하고, 옵션 정책과\n제외 재료까지 반영한 24개 테이블 데이터 모델을 설계했습니다.",
        "문서·Kiosk·Admin·Backend를 저장소로 나누고, 옵션·제외까지 담은 데이터 모델을 설계했다.",
    ),
    (
        21,
        "키오스크 주문 생성과 결제 승인 API가 미완성이고, 관리자 화면 다수가 아직\nmock 연결 단계입니다. QA 테스트 케이스도 실행 증거가 없어 검증이 남아 있습니다.",
        "주문 생성·결제 승인 E2E와 Admin mock 제거·실연동을 끝까지 검증해야 한다.",
    ),
    (
        21,
        "화면 · API · DB를 따로 확인하는 것만으로는 부족했습니다. 강사 피드백대로\nController → Service → Mapper → DB를 한 경로로 완성하는 순서가 중요했습니다.",
        "화면·API·DB는 각각 통과해도, 같은 주문번호로 한 줄로 이을 때 진짜 검증이 된다.",
    ),
]

NOTES: dict[int, str] = {
    1: "표지. ASAK = A Salad A Kiosk. 팀 다비치 · 이하진·김나연 · 멘토 남상규. AI 추천·무인결제 완성 표현 금지.",
    2: "목차. Why → Who → How → Build → Review 다섯 단계로 발표한다.",
    3: "개요. 셀프오더 트렌드 중 ‘정확한 커스터마이징’에 집중. AI 추천·성과 수치 금지.",
    4: "팀. 역할은 화면·API 단위. 통합·발표는 공동.",
    5: "절차. 문서 선행 후 화면·서버 연결. 실서버 연동·통합 검증은 진행 중.",
    11: "화면설계. 단계 분리 + 선택 유지가 UX 핵심.",
    13: "핵심 차별점. 커스터마이징 UX = 실수 없이 조합 완성. 서버 재검증 강조.",
    14: "시연 범위 B: 결제 화면 이동까지. 승인 완료는 말하지 않는다.",
    15: "운영 연동. 키오스크 주문이 Admin까지 와야 루프가 닫힌다.",
    21: "자체 평가. mock 1차 연결 ≠ DONE. 통합 검증 진행 중을 정직하게.",
    22: "마무리. ASAK는 옵션이 많은 샐러드 주문을 실수 없이 끝내고 운영 화면까지 연결한다. 질문 받겠습니다.",
}


def set_shape_text(shape, new_text: str) -> None:
    """첫 런 서식을 유지한 채 텍스트를 교체한다."""
    tf = shape.text_frame
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        return

    # 첫 단락 첫 런의 서식 백업
    first_run = None
    for p in paragraphs:
        if p.runs:
            first_run = p.runs[0]
            break

    font_snapshot = None
    if first_run is not None:
        f = first_run.font
        font_snapshot = {
            "name": f.name,
            "size": f.size,
            "bold": f.bold,
            "italic": f.italic,
            "underline": f.underline,
        }
        try:
            font_snapshot["color_rgb"] = f.color.rgb
        except Exception:
            font_snapshot["color_rgb"] = None

    # 기존 단락/런 비우기
    for i, p in enumerate(paragraphs):
        if i == 0:
            if p.runs:
                p.runs[0].text = ""
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = ""
        else:
            # 추가 단락은 비움
            for r in p.runs:
                r.text = ""
            p.text = ""

    # \n / \x0b 로 단락 분리
    parts = new_text.replace("\x0b", "\n").split("\n")
    # 첫 단락
    p0 = paragraphs[0]
    if p0.runs:
        p0.runs[0].text = parts[0]
        run = p0.runs[0]
    else:
        run = p0.add_run()
        run.text = parts[0]

    if font_snapshot:
        _apply_font(run, font_snapshot)

    # 나머지 단락: 기존 단락 재사용 또는 새로 추가
    for i, part in enumerate(parts[1:], start=1):
        if i < len(paragraphs):
            p = paragraphs[i]
            if p.runs:
                p.runs[0].text = part
                for r in p.runs[1:]:
                    r.text = ""
                r0 = p.runs[0]
            else:
                r0 = p.add_run()
                r0.text = part
        else:
            p = tf.add_paragraph()
            # 정렬 등은 첫 단락 복사 시도
            p.alignment = paragraphs[0].alignment
            r0 = p.add_run()
            r0.text = part
        if font_snapshot:
            _apply_font(r0, font_snapshot)


def _apply_font(run, snap: dict) -> None:
    f = run.font
    if snap.get("name"):
        f.name = snap["name"]
    if snap.get("size") is not None:
        f.size = snap["size"]
    if snap.get("bold") is not None:
        f.bold = snap["bold"]
    if snap.get("italic") is not None:
        f.italic = snap["italic"]
    if snap.get("underline") is not None:
        f.underline = snap["underline"]
    if snap.get("color_rgb") is not None:
        try:
            f.color.rgb = snap["color_rgb"]
        except Exception:
            pass


def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(sh.shapes)


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"source missing: {SRC}")

    shutil.copy2(SRC, DST)
    prs = Presentation(str(DST))

    # 슬라이드별 old→new 맵
    by_slide: dict[int, list[tuple[str, str]]] = {}
    for slide_no, old, new in REPLACEMENTS:
        by_slide.setdefault(slide_no, []).append((old, new))

    hit = 0
    miss: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        pairs = by_slide.get(idx, [])
        if not pairs:
            if idx in NOTES:
                slide.notes_slide.notes_text_frame.text = NOTES[idx]
            continue

        # 텍스트 인덱스
        shapes = [sh for sh in iter_shapes(slide.shapes) if sh.has_text_frame]
        used = set()
        for old, new in pairs:
            found = False
            for si, sh in enumerate(shapes):
                if si in used:
                    continue
                cur = sh.text_frame.text
                if cur == old:
                    set_shape_text(sh, new)
                    used.add(si)
                    hit += 1
                    found = True
                    break
            if not found:
                # 느슨 매칭: strip / \x0b↔\n
                old_norm = old.replace("\x0b", "\n").rstrip("\n")
                for si, sh in enumerate(shapes):
                    if si in used:
                        continue
                    cur_norm = sh.text_frame.text.replace("\x0b", "\n").rstrip("\n")
                    if cur_norm == old_norm:
                        set_shape_text(sh, new)
                        used.add(si)
                        hit += 1
                        found = True
                        break
            if not found:
                miss.append(f"S{idx}: {old[:60]!r}")

        if idx in NOTES:
            try:
                slide.notes_slide.notes_text_frame.text = NOTES[idx]
            except Exception:
                pass

    prs.save(str(DST))
    print(f"saved: {DST}")
    print(f"replacements hit={hit} miss={len(miss)}")
    for m in miss:
        print("MISS", m)


if __name__ == "__main__":
    main()
