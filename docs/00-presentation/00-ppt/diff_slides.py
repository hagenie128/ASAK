"""슬라이드별 텍스트 diff (원본 vs 현재)."""

from pathlib import Path
from pptx import Presentation

BASE = Path(__file__).resolve().parent
paths = {
    "ORIG": BASE / "ASAK_샐러드_스마트키오스크_2026_0902.broken.pptx",
    "CURR": BASE / "ASAK_샐러드_스마트키오스크_2026_0902.pptx",
}


def texts(slide):
    rows = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                rows.append((sh.top, sh.left, t[:100].replace("\n", "|")))
    return sorted(rows)


for idx in range(28):
    print(f"\n--- SLIDE {idx+1} ---")
    for label, path in paths.items():
        s = Presentation(str(path)).slides[idx]
        ts = texts(s)
        print(f"{label} ({len(ts)} shapes):")
        for top, left, t in ts[:6]:
            print(f"  {top:7} {left:7} {t}")
        if len(ts) > 6:
            print(f"  ... +{len(ts)-6} more")
