"""
Kiosk PPT 최종 빌드:
- 손상된 슬라이드(장바구니, Admin 메뉴/품절/매출)를 수정4 백업에서 복구
- Kiosk 신규 4페이지 및 캡처 반영
"""

from __future__ import annotations

import copy
import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = Path(__file__).resolve().parent
PPT_PATH = BASE / "ASAK_샐러드_스마트키오스크_2026_0902.pptx"
DONOR_PATH = BASE / "ASAK_샐러드_스마트키오스크_수정4_2026_0821.pptx"
SHOT_DIR = BASE.parent / "02-kiosk_screenshot"

IMG_2 = [(5650991, 1233811, 2263073, 4032385), (8337593, 1233812, 2287735, 4076328)]
IMG_4_ROW = [
    (438912, 1280160, 2700000, 3800000),
    (3290000, 1280160, 2700000, 3800000),
    (6141088, 1280160, 2700000, 3800000),
    (8992176, 1280160, 2700000, 3800000),
]
IMG_4_GRID = [
    (5063112, 1180000, 2900000, 2100000),
    (8200000, 1180000, 2900000, 2100000),
    (5063112, 3500000, 2900000, 2100000),
    (8200000, 3500000, 2900000, 2100000),
]

KIOSK_NEW = {
    "soldout": {
        "title": "메뉴 및 옵션 품절 처리",
        "body": (
            "메뉴와 옵션의 품절 상태를\n"
            "주문 화면에서 각각 구분하여 표시합니다.\n\n"
            "품절된 메뉴는 주문을 제한하고,\n"
            "일부 옵션만 품절된 경우에는\n"
            "해당 옵션만 선택할 수 없도록 처리했습니다."
        ),
        "highlight": (
            "품절 범위에 따라 주문 가능 여부를 다르게 제어하여 "
            "불필요한 주문 실패를 방지했습니다."
        ),
        "caption": "메뉴 전체 품절 · 개별 옵션 품절 상태",
        "images": ["10_메뉴품절상태.jpg", "11_옵션품절상태.jpg"],
        "layout": "2",
    },
    "payment": {
        "title": "결제 프로세스 및 결제수단",
        "body": (
            "주문 내역을 최종 확인한 뒤\n"
            "결제수단을 선택하고 결제를 진행합니다.\n\n"
            "결제 진행 상태와 완료 결과를\n"
            "단계별 화면으로 구분하여\n"
            "사용자가 현재 상태를 확인할 수 있도록 구성했습니다."
        ),
        "highlight": (
            "주문 확인 → 결제수단 선택 → 결제 진행 → 결제 완료의 "
            "흐름을 단계별로 구성했습니다."
        ),
        "caption": "주문 내역 확인 · 결제수단 선택 · 결제 진행 · 결제 완료",
        "footnote": "현재 결제는 실제 PG가 아닌 가상 결제 프로세스로 구현했습니다.",
        "images": ["06_주문내역확인.jpg", "07_결제수단선택.jpg", "08_결제진행.jpg", "09_결제완료.jpg"],
        "layout": "4row",
    },
    "exception": {
        "title": "예외 상황 및 사용자 상태 처리",
        "body": (
            "키오스크 사용 중 발생할 수 있는\n"
            "입력 누락 · 장바구니 상태 · 사용자 미응답 상황을\n"
            "각각 구분하여 안내합니다.\n\n"
            "단순 오류 메시지가 아니라\n"
            "사용자가 다음 행동을 알 수 있도록\n"
            "상황별 안내 팝업을 구성했습니다."
        ),
        "highlight": (
            "사용자의 실수를 사전에 방지하고, "
            "중단된 주문 흐름을 자연스럽게 이어갈 수 있도록 처리했습니다."
        ),
        "caption": "타임아웃 · 장바구니 비우기 · 빈 장바구니 · 결제수단 미선택",
        "images": [
            "12_타임아웃팝업.jpg",
            "13_장바구니_비우기_팝업.jpg",
            "14_빈장바구니_상태_팝업.jpg",
            "15_결제수단미선택_팝업.jpg",
        ],
        "layout": "4grid",
    },
    "receipt": {
        "title": "주문 완료 및 영수증 출력",
        "body": (
            "결제가 완료되면 주문 완료 상태를 안내하고\n"
            "영수증 출력 과정을 사용자에게 표시합니다.\n\n"
            "주문 종료 이후에도\n"
            "출력 상태를 확인할 수 있도록\n"
            "상태 안내 UI를 구성했습니다."
        ),
        "highlight": (
            "결제 완료 → 주문 완료 안내 → 영수증 출력까지 "
            "주문 종료 흐름을 구성했습니다."
        ),
        "caption": "결제 완료 · 영수증 출력 상태 안내",
        "footnote": (
            "RTOS Simulator와 연동하여 영수증 출력 이벤트를 처리하고, "
            "출력 상태를 Kiosk 화면에 반영했습니다."
        ),
        "images": ["09_결제완료.jpg", "17_영수증출력_상태_토스트.jpg"],
        "layout": "2",
    },
}

ADMIN_EXTRA = {
    "관리자 품절 관리": {
        "body": (
            "메뉴와 재료의 품절 상태를 \n"
            "관리자 화면에서 확인하고 관리합니다.\n\n"
            "품절 상태가 변경되면 Kiosk 주문 화면에도 \n"
            "동일한 기준을 적용해 주문할 수 없는 메뉴나 \n"
            "옵션의 선택을 제한합니다."
        ),
        "highlight": "Admin 품절 관리 → 서버 데이터 반영 → \nKiosk 주문 제한",
        "caption": "메뉴 · 재료 품절 조회 및 상태 관리",
        "footnote": "운영 중 발생하는 품절 상황을 빠르게 반영하여 고객이 주문 단계에서 품절 상품을 선택하는 오류를 줄였습니다",
    },
    "관리자 매출 관리": {
        "body": (
            "주문 데이터를 기반으로\n"
            "일별 · 월별 · 시간대별 매출 현황을 조회합니다.\n\n"
            "기간별 매출과 주문 흐름을 시각화하여\n"
            "운영자가 매장의 주요 매출 추이를\n"
            "한눈에 확인할 수 있도록 구성했습니다."
        ),
        "highlight": "주문 데이터 → 매출 집계 → 운영 지표 확인",
        "caption": "매출 현황 · 기간별 통계 · 주문 추이",
        "footnote": "주문 데이터가 단순히 저장되는 데 그치지 않고 운영에 필요한 정보로 활용될 수 있도록 관리자 화면에 시각화 했습니다.",
    },
}

EXISTING_KIOSK_IMAGES = {
    11: ["01_메뉴리스트.jpg", "03_메뉴선택후_메뉴리스트.jpg"],
    12: ["02_메뉴디테일.jpg", "03_메뉴선택후_메뉴리스트.jpg"],
    13: ["04_장바구니.jpg", "06_주문내역확인.jpg"],
}
FLOW_IMAGES = {10: ["01_메뉴리스트.jpg", "02_메뉴디테일.jpg", "04_장바구니.jpg"]}


def duplicate_slide(prs: Presentation, index: int):
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    return new_slide


def reorder_slides(prs: Presentation, order: list[int]) -> None:
    sld_id_lst = prs.slides._sldIdLst
    elements = list(sld_id_lst)
    for el in list(sld_id_lst):
        sld_id_lst.remove(el)
    for idx in order:
        sld_id_lst.append(elements[idx])


def slide_title(slide) -> str:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if shape.top < 900000 and len(text) > 4:
            if text not in ("핵심 기능 및 구현 화면", "04  수행 경과") and "ASAK" not in text:
                return text.split("\n")[0]
    return ""


def set_multiline_text(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line


def find_title_shape(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if shape.top < 900000 and len(text) > 4 and "핵심 기능" not in text and "ASAK" not in text:
            return shape
    return None


def find_body_shape(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if 1400000 < shape.top < 2800000 and len(shape.text_frame.text.strip()) > 30:
            return shape
    return None


def find_highlight_shape(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if 3200000 < shape.top < 4300000 and len(shape.text_frame.text.strip()) > 10:
            return shape
    return None


def find_caption_shape(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if 5200000 < shape.top < 5900000 and "·" in text:
            return shape
    return None


def find_footnote_shape(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if len(text) > 20 and (
            "가상 결제" in text or "PG" in text or "시연" in text or "운영" in text or "주문 데이터가" in text
        ):
            return shape
    return None


def find_page_number_shape(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if shape.left > 10000000 and shape.top > 6000000 and re.fullmatch(r"\d+", text):
            return shape
    return None


def remove_pictures(slide) -> None:
    for shape in [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]:
        shape._element.getparent().remove(shape._element)


def add_images(slide, filenames: list[str], layout: str) -> None:
    positions = {"2": IMG_2, "4row": IMG_4_ROW, "4grid": IMG_4_GRID}[layout]
    for (left, top, width, height), name in zip(positions, filenames):
        slide.shapes.add_picture(str(SHOT_DIR / name), left, top, width, height)


def replace_images_on_slide(slide, filenames: list[str]) -> None:
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if len(pics) == len(filenames):
        for pic, name in zip(pics, filenames):
            left, top, width, height = pic.left, pic.top, pic.width, pic.height
            pic._element.getparent().remove(pic._element)
            slide.shapes.add_picture(str(SHOT_DIR / name), left, top, width, height)
    else:
        remove_pictures(slide)
        layout = "2" if len(filenames) == 2 else "4row" if len(filenames) == 4 else "2"
        add_images(slide, filenames, layout)


def apply_kiosk_content(slide, content: dict) -> None:
    title = find_title_shape(slide)
    if title:
        set_multiline_text(title, content["title"])
    body = find_body_shape(slide)
    if body:
        set_multiline_text(body, content["body"])
    highlight = find_highlight_shape(slide)
    if highlight:
        set_multiline_text(highlight, content["highlight"])
    caption = find_caption_shape(slide)
    if caption:
        set_multiline_text(caption, content["caption"])
    footnote = content.get("footnote")
    if footnote:
        foot = find_footnote_shape(slide)
        if foot:
            set_multiline_text(foot, footnote)
    remove_pictures(slide)
    add_images(slide, content["images"], content["layout"])


def apply_admin_content(slide, title: str, content: dict) -> None:
    title_shape = find_title_shape(slide)
    if title_shape:
        set_multiline_text(title_shape, title)
    body = find_body_shape(slide)
    if body:
        set_multiline_text(body, content["body"])
    highlight = find_highlight_shape(slide)
    if highlight:
        set_multiline_text(highlight, content["highlight"])
    caption = find_caption_shape(slide)
    if caption:
        set_multiline_text(caption, content["caption"])
    foot = find_footnote_shape(slide)
    if foot:
        set_multiline_text(foot, content["footnote"])


def set_page_number(slide, number: int) -> None:
    shape = find_page_number_shape(slide)
    if shape:
        set_multiline_text(shape, str(number))


def import_slide_from(source_prs: Presentation, source_idx: int, target_prs: Presentation):
    source = source_prs.slides[source_idx]
    new_slide = target_prs.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    return len(target_prs.slides) - 1


def find_slide_index_by_title(prs: Presentation, title: str) -> int | None:
    for i, slide in enumerate(prs.slides):
        if slide_title(slide) == title:
            return i
    return None


def main() -> None:
    donor = Presentation(str(DONOR_PATH))
    prs = Presentation(str(PPT_PATH))

    # --- 1) 누락 슬라이드 복구 (장바구니 / Admin) ---
    cart_idx = find_slide_index_by_title(prs, "장바구니 및 결제 이동")
    admin_menu_idx = find_slide_index_by_title(prs, "관리자 메뉴 관리")
    admin_soldout_idx = find_slide_index_by_title(prs, "관리자 품절 관리")
    admin_sales_idx = find_slide_index_by_title(prs, "관리자 매출 관리")

    donor_cart = 13
    donor_menu = 15

    if cart_idx is None:
        cart_idx = import_slide_from(donor, donor_cart, prs)
        replace_images_on_slide(prs.slides[cart_idx], EXISTING_KIOSK_IMAGES[13])

    if admin_menu_idx is None:
        admin_menu_idx = import_slide_from(donor, donor_menu, prs)

    if admin_soldout_idx is None:
        admin_soldout_idx = import_slide_from(donor, donor_menu, prs)
        apply_admin_content(prs.slides[admin_soldout_idx], "관리자 품절 관리", ADMIN_EXTRA["관리자 품절 관리"])

    if admin_sales_idx is None:
        admin_sales_idx = import_slide_from(donor, donor_menu, prs)
        apply_admin_content(prs.slides[admin_sales_idx], "관리자 매출 관리", ADMIN_EXTRA["관리자 매출 관리"])

    # --- 2) Kiosk 신규 슬라이드가 없으면 추가 ---
    if find_slide_index_by_title(prs, "메뉴 및 옵션 품절 처리") is None:
        template = find_slide_index_by_title(prs, "메뉴 목록 및 카테고리 선택") or 11
        new_idx = {}
        for key in ("soldout", "payment", "exception", "receipt"):
            slide = duplicate_slide(prs, template)
            apply_kiosk_content(slide, KIOSK_NEW[key])
            new_idx[key] = len(prs.slides) - 1

        cart_idx = find_slide_index_by_title(prs, "장바구니 및 결제 이동")
        order = (
            list(range(11))
            + [11, 12, new_idx["soldout"], cart_idx, new_idx["payment"], new_idx["exception"], new_idx["receipt"]]
            + [i for i in range(14, len(prs.slides)) if i not in new_idx.values() and i != cart_idx]
        )
        reorder_slides(prs, order)

    # --- 3) 순서 재구성 ---
    prefix = list(range(11))
    kiosk_order = [
        "메뉴 목록 및 카테고리 선택",
        "메뉴 상세 및 옵션 선택",
        "메뉴 및 옵션 품절 처리",
        "장바구니 및 결제 이동",
        "결제 프로세스 및 결제수단",
        "예외 상황 및 사용자 상태 처리",
        "주문 완료 및 영수증 출력",
    ]
    admin_order = [
        "관리자 주문 관리",
        "관리자 메뉴 관리",
        "관리자 품절 관리",
        "관리자 매출 관리",
    ]
    tail_titles = [
        "수행 프로세스",
        "피드백 및 반영 내용",
        "결과물 첨부 자료 — 화면 결과",
        "결과물 첨부 자료 — 설계 · 검증 산출물",
        "자체 평가 의견",
    ]

    titles = [slide_title(s) for s in prs.slides]

    def first_index(wanted: str) -> int:
        for i, t in enumerate(titles):
            if t == wanted:
                return i
        raise RuntimeError(f"slide not found: {wanted}")

    kiosk_indices = [first_index(t) for t in kiosk_order]
    admin_indices = [first_index(t) for t in admin_order]
    tail_indices = [first_index(t) for t in tail_titles]

    closing_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "감사합니다" in shape.text_frame.text:
                closing_idx = i
                break

    order = prefix + kiosk_indices + admin_indices + tail_indices
    if closing_idx is not None and closing_idx not in order:
        order.append(closing_idx)
    reorder_slides(prs, order)

    titles = [slide_title(s) for s in prs.slides]

    def first_index(wanted: str) -> int:
        for i, t in enumerate(titles):
            if t == wanted:
                return i
        raise RuntimeError(f"slide not found: {wanted}")

    # --- 4) 캡처 교체 ---
    for idx, files in EXISTING_KIOSK_IMAGES.items():
        title = kiosk_order[idx - 11]
        si = first_index(title)
        replace_images_on_slide(prs.slides[si], files)
    for idx, files in FLOW_IMAGES.items():
        replace_images_on_slide(prs.slides[idx], files)

    # Kiosk 신규 슬라이드 내용 재적용
    for key, title in zip(
        ("soldout", "payment", "exception", "receipt"),
        kiosk_order[2:],
    ):
        apply_kiosk_content(prs.slides[first_index(title)], KIOSK_NEW[key])

    # 결과물 슬라이드 Kiosk 썸네일
    results_idx = first_index("결과물 첨부 자료 — 화면 결과")
    pics = [s for s in prs.slides[results_idx].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    for pic, name in zip(pics[:2], ["01_메뉴리스트.jpg", "02_메뉴디테일.jpg"]):
        l, t, w, h = pic.left, pic.top, pic.width, pic.height
        pic._element.getparent().remove(pic._element)
        prs.slides[results_idx].shapes.add_picture(str(SHOT_DIR / name), l, t, w, h)

    for i in range(11, len(prs.slides)):
        set_page_number(prs.slides[i], i + 1)

    tmp = PPT_PATH.with_suffix(".tmp.pptx")
    prs.save(str(tmp))
    shutil.move(tmp, PPT_PATH)
    print(f"Done: {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
