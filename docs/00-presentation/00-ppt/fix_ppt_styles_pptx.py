"""
PPT 스타일 복구 (python-pptx): 템플릿 슬라이드의 텍스트 XML 서식을 복사.
"""

from __future__ import annotations

import copy
import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

BASE = Path(__file__).resolve().parent
PPT = BASE / "ASAK_샐러드_스마트키오스크_2026_0902.pptx"
SHOT = BASE.parent / "02-kiosk_screenshot"

IMG_2 = [(5650991, 1233811, 2263073, 4032385), (8337593, 1233812, 2287735, 4076328)]
IMG_4_ROW = [
    (438912, 1280160, 2700000, 3800000),
    (3290000, 1280160, 2700000, 3800000),
    (6141088, 1280160, 2700000, 3800000),
    (8992176, 1280160, 2700000, 3800000),
]
IMG_4_GRID = [
    (5063112, 1180000, 2900000, 2100000),
    (8200000, 1180000, 2900000, 2100000),
    (5063112, 3500000, 2900000, 2100000),
    (8200000, 3500000, 2900000, 2100000),
]

KIOSK_CONTENT = {
    13: {  # slide 14
        "title": "메뉴 및 옵션 품절 처리",
        "body": (
            "메뉴와 옵션의 품절 상태를\n"
            "주문 화면에서 각각 구분하여 표시합니다.\n\n"
            "품절된 메뉴는 주문을 제한하고,\n"
            "일부 옵션만 품절된 경우에는\n"
            "해당 옵션만 선택할 수 없도록 처리했습니다."
        ),
        "highlight": (
            "품절 범위에 따라 주문 가능 여부를 다르게 제어하여 "
            "불필요한 주문 실패를 방지했습니다."
        ),
        "caption": "메뉴 전체 품절 · 개별 옵션 품절 상태",
        "images": ["10_메뉴품절상태.jpg", "11_옵션품절상태.jpg"],
        "layout": "2",
    },
    15: {  # slide 16
        "title": "결제 프로세스 및 결제수단",
        "body": (
            "주문 내역을 최종 확인한 뒤\n"
            "결제수단을 선택하고 결제를 진행합니다.\n\n"
            "결제 진행 상태와 완료 결과를\n"
            "단계별 화면으로 구분하여\n"
            "사용자가 현재 상태를 확인할 수 있도록 구성했습니다."
        ),
        "highlight": (
            "주문 확인 → 결제수단 선택 → 결제 진행 → 결제 완료의 "
            "흐름을 단계별로 구성했습니다."
        ),
        "caption": "주문 내역 확인 · 결제수단 선택 · 결제 진행 · 결제 완료",
        "footnote": "현재 결제는 실제 PG가 아닌 가상 결제 프로세스로 구현했습니다.",
        "images": ["06_주문내역확인.jpg", "07_결제수단선택.jpg", "08_결제진행.jpg", "09_결제완료.jpg"],
        "layout": "4row",
        "remove_extra_body": True,
    },
    16: {
        "title": "예외 상황 및 사용자 상태 처리",
        "body": (
            "키오스크 사용 중 발생할 수 있는\n"
            "입력 누락 · 장바구니 상태 · 사용자 미응답 상황을\n"
            "각각 구분하여 안내합니다.\n\n"
            "단순 오류 메시지가 아니라\n"
            "사용자가 다음 행동을 알 수 있도록\n"
            "상황별 안내 팝업을 구성했습니다."
        ),
        "highlight": (
            "사용자의 실수를 사전에 방지하고, "
            "중단된 주문 흐름을 자연스럽게 이어갈 수 있도록 처리했습니다."
        ),
        "caption": "타임아웃 · 장바구니 비우기 · 빈 장바구니 · 결제수단 미선택",
        "images": [
            "12_타임아웃팝업.jpg",
            "13_장바구니_비우기_팝업.jpg",
            "14_빈장바구니_상태_팝업.jpg",
            "15_결제수단미선택_팝업.jpg",
        ],
        "layout": "4grid",
    },
    17: {
        "title": "주문 완료 및 영수증 출력",
        "body": (
            "결제가 완료되면 주문 완료 상태를 안내하고\n"
            "영수증 출력 과정을 사용자에게 표시합니다.\n\n"
            "주문 종료 이후에도\n"
            "출력 상태를 확인할 수 있도록\n"
            "상태 안내 UI를 구성했습니다."
        ),
        "highlight": (
            "결제 완료 → 주문 완료 안내 → 영수증 출력까지 "
            "주문 종료 흐름을 구성했습니다."
        ),
        "caption": "결제 완료 · 영수증 출력 상태 안내",
        "footnote": (
            "RTOS Simulator와 연동하여 영수증 출력 이벤트를 처리하고, "
            "출력 상태를 Kiosk 화면에 반영했습니다."
        ),
        "images": ["09_결제완료.jpg", "17_영수증출력_상태_토스트.jpg"],
        "layout": "2",
    },
}

ADMIN_CONTENT = {
    20: {
        "title": "관리자 품절 관리",
        "body": (
            "메뉴와 재료의 품절 상태를 \n"
            "관리자 화면에서 확인하고 관리합니다.\n\n"
            "품절 상태가 변경되면 Kiosk 주문 화면에도 \n"
            "동일한 기준을 적용해 주문할 수 없는 메뉴나 \n"
            "옵션의 선택을 제한합니다."
        ),
        "highlight": "Admin 품절 관리 → 서버 데이터 반영 → \nKiosk 주문 제한",
        "caption": "메뉴 · 재료 품절 조회 및 상태 관리",
        "footnote": (
            "운영 중 발생하는 품절 상황을 빠르게 반영하여 "
            "고객이 주문 단계에서 품절 상품을 선택하는 오류를 줄였습니다"
        ),
        "remove_extra_body": True,
    },
    21: {
        "title": "관리자 매출 관리",
        "body": (
            "주문 데이터를 기반으로\n"
            "일별 · 월별 · 시간대별 매출 현황을 조회합니다.\n\n"
            "기간별 매출과 주문 흐름을 시각화하여\n"
            "운영자가 매장의 주요 매출 추이를\n"
            "한눈에 확인할 수 있도록 구성했습니다."
        ),
        "highlight": "주문 데이터 → 매출 집계 → 운영 지표 확인",
        "caption": "매출 현황 · 기간별 통계 · 주문 추이",
        "footnote": (
            "주문 데이터가 단순히 저장되는 데 그치지 않고 "
            "운영에 필요한 정보로 활용될 수 있도록 관리자 화면에 시각화 했습니다."
        ),
        "remove_extra_body": True,
    },
}


def role(top: int, text: str) -> str | None:
    if top < 900000 and "핵심 기능" not in text and "ASAK" not in text and "수행 경과" not in text:
        if len(text) > 3 and not text.isdigit():
            return "title"
    if 1400000 < top < 2800000 and len(text) > 5:
        return "body"
    if 3200000 < top < 4300000 and len(text) > 5:
        return "highlight"
    if 5200000 < top < 5900000:
        return "caption"
    if 5900000 < top < 6350000 and "ASAK" not in text and "그린컴퓨터" not in text and len(text) > 10:
        return "footnote"
    if top > 6000000 and text.isdigit():
        return "page"
    return None


def shapes_by_role(slide) -> dict[str, object]:
    out: dict[str, object] = {}
    bodies = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        r = role(shape.top, text)
        if r == "body":
            bodies.append(shape)
        elif r and r not in out:
            out[r] = shape
    if bodies:
        bodies.sort(key=lambda s: (s.left, s.top))
        out["body"] = bodies[0]
        if len(bodies) > 1:
            out["body_extra"] = bodies[1:]
    return out


def apply_styled_text(target_shape, template_shape, text: str) -> None:
    lines = [ln for ln in text.split("\n") if ln != ""] if text.count("\n\n") else text.split("\n")
    # 빈 줄 유지가 필요한 본문은 전체 split 사용
    if "\n\n" in text:
        lines = text.split("\n")

    tmpl_paras = template_shape.text_frame._txBody.findall(qn("a:p"))
    if not tmpl_paras:
        return

    tx_body = target_shape.text_frame._txBody
    body_pr = tx_body.find(qn("a:bodyPr"))
    for child in list(tx_body):
        if child is not body_pr:
            tx_body.remove(child)

    ref_para = tmpl_paras[0]
    for i, line in enumerate(lines):
        if not line and line != "":
            continue
        src = tmpl_paras[min(i, len(tmpl_paras) - 1)]
        new_p = copy.deepcopy(src)
        t_nodes = new_p.findall(".//" + qn("a:t"))
        if t_nodes:
            t_nodes[0].text = line
            for node in t_nodes[1:]:
                node.text = ""
        tx_body.append(new_p)


def remove_pictures(slide) -> None:
    for shape in [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]:
        shape._element.getparent().remove(shape._element)


def add_images(slide, files: list[str], layout: str) -> None:
    pos = {"2": IMG_2, "4row": IMG_4_ROW, "4grid": IMG_4_GRID}[layout]
    for (l, t, w, h), name in zip(pos, files):
        slide.shapes.add_picture(str(SHOT / name), l, t, w, h)


def fix_slide(slide, template_roles: dict, content: dict) -> None:
    slide_roles = shapes_by_role(slide)
    mapping = {
        "title": content.get("title"),
        "body": content.get("body"),
        "highlight": content.get("highlight"),
        "caption": content.get("caption"),
        "footnote": content.get("footnote"),
    }
    for key, text in mapping.items():
        if not text or key not in template_roles or key not in slide_roles:
            continue
        apply_styled_text(slide_roles[key], template_roles[key], text)

    if content.get("remove_extra_body"):
        for extra in slide_roles.get("body_extra", []):
            extra._element.getparent().remove(extra._element)

    if "images" in content:
        remove_pictures(slide)
        add_images(slide, content["images"], content["layout"])


def set_page(slide, num: int) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if shape.left > 10000000 and shape.top > 6000000 and re.fullmatch(r"\d+", text):
            if shape.text_frame.paragraphs[0].runs:
                shape.text_frame.paragraphs[0].runs[0].text = str(num)
            else:
                shape.text_frame.paragraphs[0].text = str(num)
            return


def main() -> None:
    backup = PPT.with_suffix(".before_style_fix.pptx")
    shutil.copy2(PPT, backup)

    prs = Presentation(str(PPT))
    kiosk_tpl = shapes_by_role(prs.slides[11])
    cart_roles = shapes_by_role(prs.slides[14])
    if "footnote" in cart_roles:
        kiosk_tpl["footnote"] = cart_roles["footnote"]

    for idx, content in {**KIOSK_CONTENT, **ADMIN_CONTENT}.items():
        fix_slide(prs.slides[idx], kiosk_tpl, content)

    # 장바구니 캡처
    cart = prs.slides[14]
    pics = [s for s in cart.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    for pic, name in zip(pics[:2], ["04_장바구니.jpg", "06_주문내역확인.jpg"]):
        l, t, w, h = pic.left, pic.top, pic.width, pic.height
        pic._element.getparent().remove(pic._element)
        cart.shapes.add_picture(str(SHOT / name), l, t, w, h)

    for i in range(11, len(prs.slides)):
        set_page(prs.slides[i], i + 1)

    out = PPT.with_suffix(".fixed.pptx")
    prs.save(str(out))
    shutil.move(out, PPT)
    print(f"스타일 복구 완료: {PPT}")
    print(f"백업: {backup}")


if __name__ == "__main__":
    main()
