"""원본 vs 현재: 도형 단위 전체 diff → JSON."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE = Path(__file__).resolve().parent
ORIG = BASE / "ASAK_샐러드_스마트키오스크_2026_0902.broken.pptx"
CURR = BASE / "ASAK_샐러드_스마트키오스크_2026_0902.pptx"
OUT = BASE / "ppt_full_diff.json"


def extract(path: Path):
    prs = Presentation(str(path))
    deck = []
    for si, slide in enumerate(prs.slides):
        shapes = []
        for sh in slide.shapes:
            item = {
                "type": str(sh.shape_type).split("(")[0].strip(),
                "top": sh.top,
                "left": sh.left,
                "width": sh.width,
                "height": sh.height,
                "name": sh.name,
            }
            if sh.has_text_frame:
                item["text"] = sh.text_frame.text
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                item["picture"] = True
            shapes.append(item)
        deck.append({"slide": si + 1, "shapes": shapes})
    return deck


def key(sh):
    return (round(sh["top"] / 10000), round(sh["left"] / 10000), sh["type"])


def main():
    orig = extract(ORIG)
    curr = extract(CURR)
    report = []
    for o_slide, c_slide in zip(orig, curr):
        sn = o_slide["slide"]
        o_map = {}
        for s in o_slide["shapes"]:
            o_map.setdefault(key(s), []).append(s)
        c_map = {}
        for s in c_slide["shapes"]:
            c_map.setdefault(key(s), []).append(s)

        slide_diffs = []
        all_keys = sorted(set(o_map) | set(c_map))
        for k in all_keys:
            os_list = o_map.get(k, [])
            cs_list = c_map.get(k, [])
            if not os_list:
                slide_diffs.append({"kind": "extra_in_current", "key": k, "current": cs_list[0]})
                continue
            if not cs_list:
                slide_diffs.append({"kind": "missing_in_current", "key": k, "original": os_list[0]})
                continue
            o, c = os_list[0], cs_list[0]
            diffs = {}
            for field in ("text", "top", "left", "width", "height"):
                if field in o or field in c:
                    ov = o.get(field)
                    cv = c.get(field)
                    if ov != cv:
                        diffs[field] = {"orig": ov, "curr": cv}
            if "picture" in o or "picture" in c:
                if bool(o.get("picture")) != bool(c.get("picture")):
                    diffs["picture"] = {"orig": o.get("picture"), "curr": c.get("picture")}
            if diffs:
                slide_diffs.append({"kind": "changed", "key": k, "name": o.get("name"), "diffs": diffs})

        if slide_diffs:
            report.append({"slide": sn, "changes": len(slide_diffs), "items": slide_diffs})

    with open(OUT, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)

    print(f"slides with diffs: {len(report)} / 28")
    for r in report:
        text_changes = sum(
            1 for it in r["items"] if it["kind"] == "changed" and "text" in it.get("diffs", {})
        )
        layout_changes = sum(
            1
            for it in r["items"]
            if it["kind"] == "changed"
            and any(k in it.get("diffs", {}) for k in ("top", "left", "width", "height"))
        )
        print(
            f"  slide {r['slide']:2}: {r['changes']} shape diffs "
            f"(text {text_changes}, layout {layout_changes})"
        )


if __name__ == "__main__":
    main()
