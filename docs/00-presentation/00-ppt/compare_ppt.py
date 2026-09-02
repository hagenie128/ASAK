"""원본(broken) vs 현재(0902) 슬라이드 내용 비교."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

BASE = Path(__file__).resolve().parent
CURRENT = BASE / "ASAK_샐러드_스마트키오스크_2026_0902.pptx"
ORIGINAL = BASE / "ASAK_샐러드_스마트키오스크_2026_0902.broken.pptx"


def slide_info(path: Path):
    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides):
        title = body = ""
        texts = []
        imgs = 0
        for sh in slide.shapes:
            if sh.shape_type == 13:
                imgs += 1
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append(t[:120])
                if sh.top < 900000 and t and "핵심" not in t and "ASAK" not in t and "수행 경과" not in t:
                    if len(t) > 3 and not t.isdigit():
                        title = t.split("\n")[0][:50]
                if 1400000 < sh.top < 2800000 and len(t) > 30:
                    body = t[:80].replace("\n", " | ")
        out.append({"i": i + 1, "title": title, "body": body, "imgs": imgs, "texts": len(texts)})
    return out


def main():
    for label, path in [("ORIGINAL", ORIGINAL), ("CURRENT", CURRENT)]:
        if not path.exists():
            print(label, "MISSING", path.name)
            continue
        print(f"\n{'='*60}\n{label} ({path.name}) slides={len(Presentation(str(path)).slides)}\n{'='*60}")
        for s in slide_info(path):
            print(f"{s['i']:2} | {s['title'][:35]:35} | {s['imgs']}img | {s['body'][:55]}")

    if ORIGINAL.exists() and CURRENT.exists():
        o = slide_info(ORIGINAL)
        c = slide_info(CURRENT)
        print(f"\n{'='*60}\nDIFFERENCES\n{'='*60}")
        n = max(len(o), len(c))
        for i in range(n):
            ot = o[i]["title"] if i < len(o) else "(없음)"
            ct = c[i]["title"] if i < len(c) else "(없음)"
            if ot != ct or (i < len(o) and i < len(c) and o[i]["imgs"] != c[i]["imgs"]):
                print(f"slide {i+1}: [{ot}] -> [{ct}]  imgs {o[i]['imgs'] if i<len(o) else '-'} -> {c[i]['imgs'] if i<len(c) else '-'}")


if __name__ == "__main__":
    main()
